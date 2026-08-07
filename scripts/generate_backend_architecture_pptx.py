"""Build a polished editable PowerPoint for MTS backend architecture.

Slide 1: Full professional architecture board (PNG)
Slide 2: Icon-rich editable component map (Azure icons when available)
Slide 3: Selling notes
"""

from __future__ import annotations

import io
import shutil
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
DIAGRAMS = ROOT / "docs" / "diagrams"
OUT = DIAGRAMS / "Multilingual-Translation-Studio-Backend-Architecture.pptx"
BOARD_PNG = DIAGRAMS / "multilingual-translator-studio-backend-architecture.png"
ICONS_DIR = DIAGRAMS / "icons"
ICON_PNG_DIR = ICONS_DIR / "png"

CREAM = RGBColor(0xFB, 0xF3, 0xE0)
CREAM_BORDER = RGBColor(0xC9, 0xB8, 0x8E)
GRAY = RGBColor(0xEE, 0xEE, 0xEE)
GRAY_BORDER = RGBColor(0x99, 0x99, 0x99)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
NAVY = RGBColor(0x1F, 0x2A, 0x44)
BLACK = RGBColor(0x22, 0x22, 0x22)
LINE = RGBColor(0x44, 0x44, 0x44)

# Prefer official Azure icon filenames (matched after zip extract).
ICON_ALIASES = {
    "app-gateway": ["applicationgateway", "application-gateway", "appgateway"],
    "apim": ["apimanagement", "api-management", "apimanagementservices"],
    "openai": ["azureopenai", "openai", "cognitive-services", "cognitiveservices"],
    "doc-intel": ["documentintelligence", "formrecognizer", "ai-document", "document-intelligence"],
    "translator": ["translator", "translatordoc", "cognitive-services-translator"],
    "service-bus": ["servicebus", "service-bus"],
    "blob": ["storageaccounts", "storage-account", "blob"],
    "postgres": ["postgresql", "azuredatabasestartforpostgresql", "sql-database"],
    "redis": ["cache-redis", "redis", "azurecacheforredis"],
    "search": ["cognitive-search", "search", "aisearch"],
    "key-vault": ["keyvaults", "key-vault", "keyvault"],
    "entra": ["entra", "activedirectory", "azureactivedirectory", "microsoft-entra-id"],
    "monitor": ["monitor", "azuremonitor"],
    "aks": ["kubernetes-services", "aks", "kubernetes"],
}


def set_fill(shape, color: RGBColor) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def set_line(shape, color: RGBColor = LINE, width_pt: float = 1.0) -> None:
    shape.line.color.rgb = color
    shape.line.width = Pt(width_pt)


def add_shadow(shape) -> None:
    """Add a soft drop shadow via OOXML."""
    sp_pr = shape._element.spPr
    # Remove existing effectLst if present
    for child in list(sp_pr):
        if "effectLst" in child.tag:
            sp_pr.remove(child)
    effect_lst = sp_pr.makeelement(qn("a:effectLst"), {})
    shadow = effect_lst.makeelement(
        qn("a:outerShdw"),
        {
            "blurRad": "50800",
            "dist": "38100",
            "dir": "2700000",
            "algn": "tl",
            "rotWithShape": "0",
        },
    )
    srgb = shadow.makeelement(qn("a:srgbClr"), {"val": "000000"})
    alpha = srgb.makeelement(qn("a:alpha"), {"val": "20000"})
    srgb.append(alpha)
    shadow.append(srgb)
    effect_lst.append(shadow)
    sp_pr.append(effect_lst)


def add_label(slide, left, top, width, height, text: str, size: int = 14, bold: bool = True, color=NAVY):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return box


def add_card(
    slide,
    left,
    top,
    width,
    height,
    fill: RGBColor,
    border: RGBColor,
    title: str,
    body: str = "",
    title_size: int = 12,
    body_size: int = 9,
    icon_path: Path | None = None,
):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    set_fill(shape, fill)
    set_line(shape, border, 1.25)
    try:
        shape.adjustments[0] = 0.1
    except Exception:
        pass
    add_shadow(shape)

    text_left = left + Inches(0.12)
    if icon_path and icon_path.exists():
        icon_size = Inches(0.38)
        slide.shapes.add_picture(
            str(icon_path),
            left + Inches(0.1),
            top + Inches(0.12),
            icon_size,
            icon_size,
        )
        text_left = left + Inches(0.55)

    # text overlay box (transparent)
    tb = slide.shapes.add_textbox(text_left, top + Inches(0.08), width - (text_left - left) - Inches(0.1), height - Inches(0.12))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(title_size)
    run.font.bold = True
    run.font.color.rgb = BLACK
    run.font.name = "Calibri"
    if body:
        p2 = tf.add_paragraph()
        run2 = p2.add_run()
        run2.text = body
        run2.font.size = Pt(body_size)
        run2.font.color.rgb = BLACK
        run2.font.name = "Calibri"
    return shape


def find_icon_svg(keywords: list[str]) -> Path | None:
    azure_root = ICONS_DIR / "azure"
    if not azure_root.exists():
        return None
    svgs = list(azure_root.rglob("*.svg"))
    lowered = [(p, p.name.lower().replace(" ", "").replace("_", "").replace("-", "")) for p in svgs]
    for key in keywords:
        key_n = key.lower().replace(" ", "").replace("_", "").replace("-", "")
        for path, name in lowered:
            if key_n in name:
                return path
    return None


def svg_to_png(svg_path: Path, png_path: Path, size: int = 128) -> bool:
    """Convert SVG → PNG using svglib/reportlab when available."""
    try:
        from reportlab.graphics import renderPM
        from svglib.svglib import svg2rlg
    except ImportError:
        return False
    try:
        drawing = svg2rlg(str(svg_path))
        if drawing is None:
            return False
        # scale
        scale = size / max(drawing.width or size, drawing.height or size, 1)
        drawing.width *= scale
        drawing.height *= scale
        drawing.scale(scale, scale)
        png_path.parent.mkdir(parents=True, exist_ok=True)
        renderPM.drawToFile(drawing, str(png_path), fmt="PNG")
        return png_path.exists()
    except Exception:
        return False


def prepare_icons() -> dict[str, Path]:
    ICON_PNG_DIR.mkdir(parents=True, exist_ok=True)
    prepared: dict[str, Path] = {}
    for alias, keywords in ICON_ALIASES.items():
        out = ICON_PNG_DIR / f"{alias}.png"
        if out.exists():
            prepared[alias] = out
            continue
        svg = find_icon_svg(keywords)
        if svg and svg_to_png(svg, out):
            prepared[alias] = out
    return prepared


def add_arrow(slide, x1, y1, x2, y2) -> None:
    connector = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    connector.line.color.rgb = LINE
    connector.line.width = Pt(1.75)
    ln = connector.line._ln
    ln.append(ln.makeelement(qn("a:tailEnd"), {"type": "triangle", "w": "med", "len": "med"}))


def build() -> Path:
    icons = prepare_icons()
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ---------- Slide 1: polished board image ----------
    s1 = prs.slides.add_slide(prs.slide_layouts[6])
    # light background
    bg = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
    set_fill(bg, RGBColor(0xF7, 0xF7, 0xF5))
    bg.line.fill.background()

    add_label(
        s1,
        Inches(0.35),
        Inches(0.12),
        Inches(12),
        Inches(0.35),
        "SA - Multilingual Translation Studio (Backend)",
        18,
        True,
        NAVY,
    )
    if BOARD_PNG.exists():
        # Keep margins for title; image fills rest of slide
        s1.shapes.add_picture(
            str(BOARD_PNG),
            Inches(0.2),
            Inches(0.5),
            width=Inches(12.9),
            height=Inches(6.8),
        )
    else:
        add_label(s1, Inches(1), Inches(3), Inches(10), Inches(1), "Architecture PNG missing", 16, True)

    # ---------- Slide 2: icon-rich editable map ----------
    s2 = prs.slides.add_slide(prs.slide_layouts[6])
    bg2 = s2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
    set_fill(bg2, RGBColor(0xF7, 0xF7, 0xF5))
    bg2.line.fill.background()

    add_label(s2, Inches(0.3), Inches(0.12), Inches(10), Inches(0.35), "Editable component map (drag shapes / replace icons)", 16, True, NAVY)

    # legend
    leg_c = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.7), Inches(0.15), Inches(0.28), Inches(0.2))
    set_fill(leg_c, CREAM)
    set_line(leg_c, CREAM_BORDER)
    add_label(s2, Inches(10.05), Inches(0.12), Inches(3), Inches(0.25), "New Application", 10, False, BLACK)
    leg_g = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.7), Inches(0.4), Inches(0.28), Inches(0.2))
    set_fill(leg_g, GRAY)
    set_line(leg_g, GRAY_BORDER)
    add_label(s2, Inches(10.05), Inches(0.37), Inches(3), Inches(0.25), "Shared Platform", 10, False, BLACK)

    add_card(s2, Inches(0.3), Inches(0.85), Inches(2.2), Inches(1.35), GRAY, GRAY_BORDER, "API Consumers", "Partner · BFF · Mobile · External", 12, 9)
    add_card(s2, Inches(0.45), Inches(2.45), Inches(1.9), Inches(0.7), GRAY, GRAY_BORDER, "F5", "Load balancer", 12, 9)
    add_card(
        s2,
        Inches(0.35),
        Inches(3.4),
        Inches(2.1),
        Inches(0.85),
        GRAY,
        GRAY_BORDER,
        "App Gateway",
        "TLS · WAF · routing",
        12,
        9,
        icons.get("app-gateway"),
    )

    # CAT3 container
    env = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.75), Inches(0.8), Inches(5.5), Inches(4.7))
    set_fill(env, WHITE)
    set_line(env, RGBColor(0x88, 0x88, 0x88), 1.75)
    try:
        env.adjustments[0] = 0.04
    except Exception:
        pass
    add_shadow(env)
    add_label(s2, Inches(2.95), Inches(0.9), Inches(5), Inches(0.3), "MTS Backend Environment (CAT 3)", 13, True, NAVY)

    add_card(
        s2,
        Inches(2.95),
        Inches(1.25),
        Inches(5.1),
        Inches(1.05),
        CREAM,
        CREAM_BORDER,
        "NS1 · API Gateway Services",
        "FastAPI /api/v1 · Upload · Status · Pages · Cancel/Retry · Downloads · JWT",
        12,
        9,
        icons.get("aks"),
    )
    add_card(
        s2,
        Inches(2.95),
        Inches(2.45),
        Inches(5.1),
        Inches(1.55),
        CREAM,
        CREAM_BORDER,
        "NS2 · Orchestrator & Processing",
        "Intake/Extract → Gateway/Profiles → Translate → Export\nData Security Gateway · Rule/Profile library · JWT between services",
        12,
        9,
        icons.get("aks"),
    )
    add_card(
        s2,
        Inches(2.95),
        Inches(4.15),
        Inches(5.1),
        Inches(1.15),
        CREAM,
        CREAM_BORDER,
        "NS3 · Workers & Agents",
        "Page-range DI worker · Translation worker · Retention · Prompt library · Token service",
        12,
        9,
        icons.get("aks"),
    )

    add_card(s2, Inches(8.55), Inches(0.85), Inches(1.5), Inches(0.7), GRAY, GRAY_BORDER, "APIM", "JWT · throttle", 11, 8, icons.get("apim"))
    add_card(s2, Inches(10.25), Inches(0.85), Inches(1.6), Inches(0.7), GRAY, GRAY_BORDER, "App Gateway", "AI path", 11, 8, icons.get("app-gateway"))

    ai = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.55), Inches(1.7), Inches(4.45), Inches(2.85))
    set_fill(ai, CREAM)
    set_line(ai, CREAM_BORDER, 1.75)
    try:
        ai.adjustments[0] = 0.05
    except Exception:
        pass
    add_shadow(ai)
    add_label(s2, Inches(8.7), Inches(1.8), Inches(4), Inches(0.28), "AI HUB (CAT 1)", 13, True, NAVY)
    add_card(s2, Inches(8.75), Inches(2.15), Inches(4.05), Inches(0.65), WHITE, CREAM_BORDER, "Azure OpenAI", "GPT structured translation", 11, 8, icons.get("openai"))
    add_card(s2, Inches(8.75), Inches(2.9), Inches(4.05), Inches(0.65), WHITE, CREAM_BORDER, "Document Intelligence 4.0", "prebuilt-layout · OCR · languages", 11, 8, icons.get("doc-intel"))
    add_card(s2, Inches(8.75), Inches(3.65), Inches(4.05), Inches(0.65), WHITE, CREAM_BORDER, "Azure AI Translator", "Non-LLM route", 11, 8, icons.get("translator"))

    add_card(s2, Inches(8.55), Inches(4.75), Inches(2.1), Inches(0.7), GRAY, GRAY_BORDER, "AZ Entra ID", "App roles · SPN", 11, 8, icons.get("entra"))
    add_card(s2, Inches(10.85), Inches(4.75), Inches(2.15), Inches(0.7), GRAY, GRAY_BORDER, "Identity Governance", "Enterprise IDP", 11, 8)

    # Data layer
    data = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.75), Inches(5.7), Inches(10.25), Inches(1.15))
    set_fill(data, CREAM)
    set_line(data, CREAM_BORDER, 1.5)
    add_shadow(data)
    add_label(s2, Inches(2.9), Inches(5.75), Inches(9.8), Inches(0.25), "Data Storage Layer · CMK encryption", 11, True, NAVY)
    stores = [
        (2.9, "Service Bus", "service-bus"),
        (4.85, "Blob Storage", "blob"),
        (6.8, "Postgres DB", "postgres"),
        (8.75, "Redis Cache", "redis"),
        (10.7, "AI Search", "search"),
    ]
    for x, title, key in stores:
        add_card(s2, Inches(x), Inches(6.05), Inches(1.8), Inches(0.7), WHITE, CREAM_BORDER, title, "", 10, 8, icons.get(key))

    add_arrow(s2, Inches(2.5), Inches(3.8), Inches(2.75), Inches(3.8))
    add_arrow(s2, Inches(8.25), Inches(3.0), Inches(8.55), Inches(1.2))

    # Platform strip
    platforms = [
        (0.3, "Key Vault", "key-vault"),
        (2.2, "Entra Identity", "entra"),
        (4.2, "Azure Monitor", "monitor"),
        (6.2, "Log Analytics", None),
        (8.1, "GitLab CI/CD", None),
        (10.0, "Artifacts", None),
        (11.7, "Private EP", None),
    ]
    for x, title, key in platforms:
        add_card(s2, Inches(x), Inches(7.05), Inches(1.55), Inches(0.35), GRAY, GRAY_BORDER, title, "", 9, 7, icons.get(key) if key else None)

    # ---------- Slide 3: notes ----------
    s3 = prs.slides.add_slide(prs.slide_layouts[6])
    bg3 = s3.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
    set_fill(bg3, WHITE)
    bg3.line.fill.background()
    add_label(s3, Inches(0.4), Inches(0.3), Inches(12), Inches(0.4), "How to use / edit this deck", 18, True, NAVY)
    notes = [
        "Slide 1 — Presentation board: high-fidelity architecture image for customer reviews.",
        "Slide 2 — Editable map: native PowerPoint shapes + Azure icons (drag, recolor, rename).",
        "Icons source: Microsoft Azure Architecture Center official SVG pack (when downloaded into docs/diagrams/icons).",
        "",
        "Sell as backend-only: customers bring UI; they integrate via App Gateway / APIM.",
        "CAT 3 = MTS namespaces (API · Orchestrator · Workers).  CAT 1 = AI Hub via APIM.",
        "Data plane: Service Bus · Blob · Postgres · Redis · AI Search with CMK.",
        "",
        "Also editable in diagrams.net: Multilingual-Translation-Studio-Backend-Architecture.drawio",
        "Regenerate: backend\\.venv\\Scripts\\python.exe scripts\\generate_backend_architecture_pptx.py",
    ]
    tb = s3.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(12.2), Inches(5.8))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(notes):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = line
        run.font.size = Pt(15)
        run.font.color.rgb = BLACK
        run.font.name = "Calibri"
        run.font.bold = line.startswith("Slide") or line.startswith("Sell") or line.startswith("Also")

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    return OUT


if __name__ == "__main__":
    print(build())
    print("icons prepared:", list(prepare_icons().keys()))
